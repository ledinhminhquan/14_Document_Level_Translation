"""Generate the submission slides.pptx (python-pptx) — ~12 concise slides for the doctrans
document-translation system. Degrades to a Markdown outline if python-pptx is unavailable.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from ..config import AppConfig, run_dir
from ..logging_utils import get_logger
from . import charts as charts_mod
from .artifact_loader import doc_metric, load_artifacts, model_version, mt_metric

logger = get_logger(__name__)


def _slides(cfg: AppConfig, arts: Dict[str, Any]) -> List[Tuple[str, List[str]]]:
    mc = mt_metric(arts, "model", "chrf")
    dc = mt_metric(arts, "dictionary", "chrf")
    dchrf = doc_metric(arts, "doc_chrf")
    dsps = doc_metric(arts, "doc_sps")
    res = (f"MT chrF {mc:.1f} vs dictionary {dc:.1f}" if (mc is not None and dc is not None)
           else "train + evaluate to populate results")
    docline = (f"document chrF {dchrf:.1f} at structure-preservation {dsps:.2f}"
               if (dchrf is not None and dsps is not None) else "document chrF + structure-preservation (SPS)")
    d = f"{cfg.mt.src_lang} -> {cfg.mt.tgt_lang}"
    return [
        ("Structured Document-Level Translation",
         [f"{cfg.author} - Student {cfg.student_id}", "NLP in Industry - Final Assignment",
          f"Translate Markdown/HTML/JSON docs ({d}) preserving structure",
          "parse -> mask -> translate-with-context -> reassemble -> validate",
          "Structure-breaking output is flagged for human review"]),
        ("Business Problem & Motivation",
         ["Raw MT corrupts structure: eaten markers, translated URLs/code, dropped placeholders",
          "Localization teams need format-stable, context-consistent translations",
          "The fix = the localization spine + a trainable, context-aware MT core + a QA agent",
          "Only the MT model is trained (the NLP heart); the structure layer is deterministic"]),
        ("Proposed Solution",
         ["Parse per format -> interleaved skeleton (literal) + translatable Segments",
          "Mask non-translatable spans (code/URLs/placeholders/tags) to [[PHn]] sentinels",
          "Translate prose with document-level concat-k context; restore masks",
          "Reassemble byte-stably + validate structure (D5)"]),
        ("System Architecture",
         ["parse (D1, round-trip identity check) -> segment + mask (D2)",
          "-> translate-with-context (D3) -> verify (D4)",
          "-> reassemble + structure-validate (D5)",
          "model never sees structure; structure never sees the model"]),
        ("Data & Languages",
         ["Fine-tune: Helsinki-NLP/opus-100 en-fr (license unknown flag)",
          "Context: news_commentary en-fr (monotonic id -> real document adjacency)",
          "No structured parallel corpus exists -> structure is synthetic",
          f"Default {d}; offline seed = Markdown docs + gold + en->fr dictionary"]),
        ("The Trainable MT Core + Context",
         ["Base: facebook/m2m100_418M (MIT, 1024 positions) - reuses P13 s2st plumbing",
          "Document context = concat-k (k previous sentences + <BRK>, k-to-1)",
          "Seq2SeqTrainer, chrF selection, bf16/tf32, resume-safe",
          "Baselines: zero-shot MT + dictionary + identity floor"]),
        ("Agentic AI Component (D1-D5)",
         ["Deterministic FSM + optional LLM terminology brain (OFF by default)",
          "D1 format+parse (round-trip check) - D2 segment+mask gate",
          "D3 translate (length-ratio sanity) - D4 verify (placeholder-retention HARD + chrF)",
          "D5 reassemble + structure-validation gate (signature match + re-parse)"]),
        ("Evaluation Results",
         [res, docline,
          "chrF + Structure-Preservation Score (SPS) + placeholder-retention reported jointly",
          "Context story: discourse fixes show on contrastive sets, not raw BLEU"]),
        ("Deployment Overview",
         ["FastAPI /translate-text + /translate-document (multipart) + /healthz + /version",
          "Gradio demo (paste doc -> translated doc + structure report)",
          "Docker + HF Space; lazy deps, line/regex parser + dictionary offline fallback",
          "Metadata-only request logging"]),
        ("Continual Learning & Monitoring",
         ["Collect domain docs -> re-fine-tune MT -> promote if chrF non-regressing",
          "Structure layer is deterministic code (versioned, not trained)",
          "monitor-log: needs-review rate + mean SPS + latency drift",
          "Rising needs-review / falling SPS = drift signal"]),
        ("Ethics, Privacy & Risks",
         ["Localization aid - structure-breaking / low-confidence output flagged for review",
          "Mistranslation harm -> verify + structure gates + abstention",
          "Placeholder/markup loss -> hard retention gate + fail-soft (keep source)",
          "License hygiene -> opus/news_commentary unknown, NLLB CC-BY-NC; clean = m2m100 MIT"]),
        ("Key Takeaways & Future Work",
         ["A structure-stable, context-aware, debuggable document-translation pipeline",
          "Hard placeholder-retention + structure-validation gates make it safe",
          "Future: DOCX/PDF formats, word-aligner detag, ContraPro context eval, terminology memory",
          "Future: larger context windows + a document-level MT model"]),
    ]


def generate_slides(cfg: AppConfig, title: Optional[str] = None, author: Optional[str] = None,
                    out_path: Optional[str] = None) -> str:
    arts = load_artifacts(cfg)
    out_path = Path(out_path) if out_path else run_dir() / "report" / "slides.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    slides = _slides(cfg, arts)
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except Exception as exc:
        logger.warning("python-pptx unavailable (%s); writing markdown outline", exc)
        md = "\n\n".join(f"## {t}\n" + "\n".join(f"- {b}" for b in bs) for t, bs in slides)
        alt = out_path.with_suffix(".md")
        alt.write_text(md, encoding="utf-8")
        return str(alt)

    try:
        chart = charts_mod.fidelity_chart(arts.get("eval") or {}, run_dir() / "report" / "slide_fidelity.png")
    except Exception:
        chart = None
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    accent = RGBColor(0x2B, 0x6C, 0xB0)
    for i, (t, bullets) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
        tf = bar.text_frame; tf.text = t
        tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.5),
                                        Inches(8.3 if (i == 7 and chart) else 12), Inches(5.4))
        bt = body.text_frame; bt.word_wrap = True
        for j, bp in enumerate(bullets):
            p = bt.paragraphs[0] if j == 0 else bt.add_paragraph()
            p.text = "-  " + bp; p.font.size = Pt(20); p.space_after = Pt(10)
        if i == 7 and chart:
            slide.shapes.add_picture(str(chart), Inches(8.9), Inches(1.7), width=Inches(4.0))
        foot = slide.shapes.add_textbox(Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4))
        foot.text_frame.text = f"{title or cfg.project_title} - {author or cfg.author} ({cfg.student_id})"
        foot.text_frame.paragraphs[0].font.size = Pt(9)
    prs.save(str(out_path))
    logger.info("Slides -> %s", out_path)
    return str(out_path)


__all__ = ["generate_slides"]
